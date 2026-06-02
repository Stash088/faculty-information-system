#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import re
import html

# Read the HTML file
with open('/Users/amir/Desktop/back_work_final/presentation.html', 'r', encoding='utf-8') as f:
    html_content = f.read()

# Extract CSS colors
colors = {}
css_match = re.search(r'<style>(.*?)</style>', html_content, re.DOTALL)
if css_match:
    css = css_match.group(1)
    for name, var in [('bg', '--bg'), ('accent', '--accent'), ('text', '--text'), 
                      ('text_light', '--text-light'), ('border', '--border'),
                      ('success', '--success'), ('warning', '--warning')]:
        m = re.search(var + r':\s*(#[a-fA-F0-9]+)', css)
        if m:
            colors[name] = m.group(1)

def hex_to_rgb(h):
    h = h.lstrip('#')
    return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))

# Create presentation (16:9)
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

def add_textbox(slide, text, left, top, width, height, font_size=18, bold=False, 
                color='#222222', align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = 'Arial'
    r, g, b = hex_to_rgb(color)
    run.font.color.rgb = RGBColor(r, g, b)
    return txBox

def add_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    r, g, b = hex_to_rgb(fill_color)
    shape.fill.fore_color.rgb = RGBColor(r, g, b)
    if border_color:
        shape.line.color.rgb = RGBColor(*hex_to_rgb(border_color))
    else:
        shape.line.fill.background()
    return shape

# Parse slides using proper div counting
def find_slides(html):
    slides = []
    slide_starts = list(re.finditer(r'<div class="slide[^"]*" data-slide="(\d+)">', html))
    for i, match in enumerate(slide_starts):
        slide_num = match.group(1)
        start = match.start()
        end = slide_starts[i+1].start() if i+1 < len(slide_starts) else len(html)
        depth = 1
        pos = match.end()
        while depth > 0 and pos < end:
            next_open = html.find('<div', pos)
            next_close = html.find('</div>', pos)
            if next_close < next_open or next_open == -1:
                depth -= 1
                pos = next_close + 6
            else:
                depth += 1
                pos = next_open + 5
        slides.append((slide_num, html[start:pos]))
    return slides

slides = find_slides(html_content)
print(f"Found {len(slides)} slides")

for slide_num, content in slides:
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    add_rect(slide, 0, 0, W, H, colors.get('bg', '#ffffff'))
    content = content.strip()
    
    if 'slide-title' in content:
        # Title slide
        add_textbox(slide, 'АГУ', W/2 - Inches(0.6), Inches(0.8), Inches(1.2), Inches(0.8), 
                    font_size=28, bold=True, color=colors.get('accent', '#2c5282'), align=PP_ALIGN.CENTER)
        
        title_match = re.search(r'<h1[^>]*>(.*?)</h1>', content, re.DOTALL)
        if title_match:
            title = html.unescape(re.sub(r'<[^>]+>', '', title_match.group(1)))
            add_textbox(slide, title, Inches(0.8), Inches(1.8), Inches(11.7), Inches(1.2),
                        font_size=24, color=colors.get('text', '#222222'), align=PP_ALIGN.CENTER)
        
        sub_match = re.search(r'class="subtitle">(.*?)</p>', content)
        if sub_match:
            sub = html.unescape(sub_match.group(1))
            add_textbox(slide, sub, W/2 - Inches(4), Inches(3.0), Inches(8), Inches(0.5),
                        font_size=18, italic=True, color=colors.get('text_light', '#555555'), align=PP_ALIGN.CENTER)
        
        meta_items = re.findall(r'meta-value">(.*?)</div>', content, re.DOTALL)
        if len(meta_items) >= 2:
            add_textbox(slide, 'Студент: ' + html.unescape(meta_items[0]), 
                        W/2 - Inches(2.5), Inches(4.0), Inches(5), Inches(0.4),
                        font_size=16, color=colors.get('text', '#222222'), align=PP_ALIGN.CENTER)
            add_textbox(slide, 'Научный руководитель: ' + html.unescape(meta_items[1]),
                        W/2 - Inches(3), Inches(4.5), Inches(6), Inches(0.4),
                        font_size=16, color=colors.get('text', '#222222'), align=PP_ALIGN.CENTER)
        
        footer_match = re.search(r'class="footer">(.*?)</div>', content, re.DOTALL)
        if footer_match:
            footer = html.unescape(re.sub(r'<[^>]+>', ' ', footer_match.group(1)))
            add_textbox(slide, footer, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.8),
                        font_size=14, color=colors.get('text_light', '#555555'), align=PP_ALIGN.CENTER)
    
    elif 'slide-conclusion' in content:
        # Check if this is the "clean" conclusion slide (only h1)
        h1_match = re.search(r'<h1[^>]*>(.*?)</h1>', content)
        if h1_match and 'Спасибо за внимание' in h1_match.group(1):
            # Simple conclusion - just the heading
            add_textbox(slide, 'Спасибо за внимание', Inches(0.8), Inches(3.5), Inches(11.7), Inches(1.2),
                        font_size=44, bold=True, color=colors.get('text', '#222222'), align=PP_ALIGN.CENTER)
        else:
            # Full conclusion with check icon and metadata
            add_textbox(slide, '✓', W/2 - Inches(0.5), Inches(1.5), Inches(1), Inches(1),
                        font_size=48, color=colors.get('success', '#276749'), align=PP_ALIGN.CENTER)
            if h1_match:
                t = html.unescape(h1_match.group(1))
                add_textbox(slide, t, Inches(0.8), Inches(2.8), Inches(11.7), Inches(0.8),
                            font_size=36, bold=True, color=colors.get('text', '#222222'), align=PP_ALIGN.CENTER)
            sub_match = re.search(r'class="subtitle">(.*?)</p>', content)
            if sub_match:
                add_textbox(slide, html.unescape(sub_match.group(1)), W/2 - Inches(3), Inches(3.7), Inches(6), Inches(0.5),
                            font_size=20, italic=True, color=colors.get('text_light', '#555555'), align=PP_ALIGN.CENTER)
            add_textbox(slide, 'Кузьминова А.А.', W/2 - Inches(2), Inches(5.2), Inches(4), Inches(0.4),
                        font_size=18, color=colors.get('text', '#222222'), align=PP_ALIGN.CENTER)
            add_textbox(slide, 'доцент В.В. Бучацкая', W/2 - Inches(2.5), Inches(5.7), Inches(5), Inches(0.4),
                        font_size=16, color=colors.get('text_light', '#555555'), align=PP_ALIGN.CENTER)
    
    else:
        # Regular slide
        label_match = re.search(r'<span class="section-label">(.*?)</span>', content)
        if label_match:
            add_textbox(slide, label_match.group(1), W - Inches(2), Inches(0.25), Inches(1.8), Inches(0.4),
                        font_size=11, color=colors.get('text_light', '#555555'), align=PP_ALIGN.RIGHT)
        
        h2_match = re.search(r'<h2[^>]*>(.*?)</h2>', content)
        if h2_match:
            t = html.unescape(h2_match.group(1))
            add_textbox(slide, t, Inches(0.6), Inches(0.35), Inches(10), Inches(0.7),
                        font_size=32, bold=True, color=colors.get('accent', '#2c5282'))
        
        # Large image slide (slides 9-12 with interface screenshots)
        img_match = re.search(r'<img src="figures/([^"]+)"[^>]*alt="([^"]*)"', content)
        if img_match:
            img_name = img_match.group(1)
            img_label = img_match.group(2)
            img_path = '/Users/amir/Desktop/back_work_final/figures/' + img_name
            try:
                # Center image, 75% of slide
                img = slide.shapes.add_picture(img_path, Inches(2.0), Inches(1.3), Inches(9.3), Inches(5.2))
            except Exception as e:
                print(f"  Could not add image {img_name}: {e}")
            
            # Description below
            desc_match = re.search(r'<p[^>]*style="text-align.*?center.*?">(.*?)</p>', content)
            if desc_match:
                add_textbox(slide, html.unescape(desc_match.group(1)), Inches(1), Inches(6.8), Inches(11.3), Inches(0.5),
                            font_size=16, italic=True, color=colors.get('text_light', '#555555'), align=PP_ALIGN.CENTER)
            continue
        
        # Problem grid
        problems = re.findall(r'<h4>(.*?)</h4>.*?<p[^>]*>(.*?)</p>', content, re.DOTALL)
        if problems and len(problems) == 4:
            positions = [(0.5, 1.3), (6.7, 1.3), (0.5, 3.9), (6.7, 3.9)]
            for i, (title, desc) in enumerate(problems[:4]):
                x, y = positions[i]
                bar = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(0.1), Inches(2.1))
                bar.fill.solid()
                bar.fill.fore_color.rgb = RGBColor(*hex_to_rgb(colors.get('warning', '#c05621')))
                bar.line.fill.background()
                add_textbox(slide, html.unescape(title), Inches(x + 0.2), Inches(y), Inches(5.8), Inches(0.45),
                            font_size=17, bold=True, color=colors.get('warning', '#c05621'))
                add_textbox(slide, html.unescape(desc), Inches(x + 0.2), Inches(y + 0.45), Inches(5.8), Inches(1.6),
                            font_size=15, color=colors.get('text_light', '#555555'))
        
        # Num list - get ul first, then find all li inside
        num_list_match = re.search(r'<ul[^>]*class="num-list"[^>]*>(.*?)</ul>', content, re.DOTALL)
        if num_list_match:
            li_items = re.findall(r'<li>(.*?)</li>', num_list_match.group(1), re.DOTALL)
            for i, item in enumerate(li_items[:5]):
                text = html.unescape(re.sub(r'<[^>]+>', '', item))
                circle = slide.shapes.add_shape(9, Inches(0.5), Inches(1.2 + i*0.9), Inches(0.38), Inches(0.38))
                circle.fill.solid()
                circle.fill.fore_color.rgb = RGBColor(*hex_to_rgb(colors.get('accent', '#2c5282')))
                circle.line.fill.background()
                add_textbox(slide, str(i+1), Inches(0.5), Inches(1.2 + i*0.9), Inches(0.38), Inches(0.38),
                            font_size=13, bold=True, color='#ffffff', align=PP_ALIGN.CENTER)
                add_textbox(slide, text, Inches(1.0), Inches(1.25 + i*0.9), Inches(10), Inches(0.45),
                            font_size=17, color=colors.get('text', '#222222'))
        
        # Tech items
        tech_items = re.findall(r'tech-item.*?tech-name">(.*?)</div>.*?tech-desc">(.*?)</div>', content, re.DOTALL)
        if tech_items:
            cols = [[], []]
            for i, (name, desc) in enumerate(tech_items):
                cols[i % 2].append((html.unescape(name), html.unescape(desc)))
            for col_idx, items in enumerate(cols):
                base_x = 0.5 + col_idx * 6.3
                for i, (name, desc) in enumerate(items):
                    y = 1.3 + i * 1.0
                    badge = slide.shapes.add_shape(1, Inches(base_x), Inches(y + 0.05), Inches(0.9), Inches(0.32))
                    badge.fill.solid()
                    badge.fill.fore_color.rgb = RGBColor(*hex_to_rgb(colors.get('accent', '#2c5282')))
                    badge.line.fill.background()
                    add_textbox(slide, name, Inches(base_x + 1.1), Inches(y), Inches(4.5), Inches(0.4),
                                font_size=17, bold=True, color=colors.get('text', '#222222'))
                    add_textbox(slide, desc, Inches(base_x + 1.1), Inches(y + 0.38), Inches(4.5), Inches(0.4),
                                font_size=13, color=colors.get('text_light', '#555555'))
        
        # Cards (security slide) - use h3 with optional style attribute
        card_divs = re.findall(r'<div class="card">(.*?)</div>', content, re.DOTALL)
        if card_divs and len(card_divs) >= 3:
            for i, card_content in enumerate(card_divs[:3]):
                x = 0.5 + i * 4.15
                h3_match = re.search(r'<h3[^>]*>(.*?)</h3>', card_content)
                p_match = re.search(r'<p>(.*?)</p>', card_content)
                if h3_match and p_match:
                    title = html.unescape(h3_match.group(1))
                    desc = html.unescape(p_match.group(1))
                    add_rect(slide, Inches(x), Inches(1.2), Inches(3.8), Inches(1.9), '#f5f5f5', colors.get('border', '#dddddd'))
                    add_textbox(slide, title, Inches(x + 0.15), Inches(1.3), Inches(3.5), Inches(0.45),
                                font_size=16, bold=True, color=colors.get('accent', '#2c5282'))
                    add_textbox(slide, desc, Inches(x + 0.15), Inches(1.8), Inches(3.5), Inches(1.2),
                                font_size=13, color=colors.get('text_light', '#555555'))
        
        # Stats
        stat_values = re.findall(r'stat-value">(.*?)</div>', content)
        if stat_values and len(stat_values) >= 4:
            for i, val in enumerate(stat_values[:4]):
                x = 0.5 + i * 3.1
                add_rect(slide, Inches(x), Inches(3.4), Inches(2.8), Inches(1.1), '#f5f5f5', colors.get('border', '#dddddd'))
                add_textbox(slide, html.unescape(val), Inches(x), Inches(3.5), Inches(2.8), Inches(0.55),
                            font_size=24, bold=True, color=colors.get('accent', '#2c5282'), align=PP_ALIGN.CENTER)
        
        # Data table (roles)
        table_rows = re.findall(r'<td[^>]*>(.*?)</td>', content, re.DOTALL)
        if table_rows and len(table_rows) >= 8:
            rows_data = []
            for i in range(4):
                if i*2+1 < len(table_rows):
                    rows_data.append((html.unescape(table_rows[i*2]), html.unescape(table_rows[i*2+1])))
            
            add_rect(slide, Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.5), colors.get('accent', '#2c5282'))
            add_textbox(slide, 'Роль', Inches(0.6), Inches(1.25), Inches(1.7), Inches(0.4),
                        font_size=15, bold=True, color='#ffffff')
            add_textbox(slide, 'Функционал', Inches(2.4), Inches(1.25), Inches(10), Inches(0.4),
                        font_size=15, bold=True, color='#ffffff')
            
            for ri, (role, func) in enumerate(rows_data):
                y = 1.8 + ri * 0.7
                bg = '#f5f5f5' if ri % 2 == 0 else '#ffffff'
                add_rect(slide, Inches(0.5), Inches(y), Inches(12.3), Inches(0.65), bg)
                add_textbox(slide, role, Inches(0.6), Inches(y + 0.08), Inches(1.7), Inches(0.5),
                            font_size=14, bold=True, color=colors.get('accent', '#2c5282'))
                add_textbox(slide, func, Inches(2.4), Inches(y + 0.08), Inches(10), Inches(0.5),
                            font_size=14, color=colors.get('text', '#222222'))
        
        # Result cards
        result_cards = re.findall(r'result-card.*?<h4>(.*?)</h4>.*?<p>(.*?)</p>', content, re.DOTALL)
        if result_cards and len(result_cards) == 3:
            for i, (title, desc) in enumerate(result_cards):
                x = 0.5 + i * 4.15
                add_rect(slide, Inches(x), Inches(1.2), Inches(3.8), Inches(1.9), '#f0fff4', colors.get('success', '#276749'))
                add_textbox(slide, html.unescape(title), Inches(x + 0.15), Inches(1.3), Inches(3.5), Inches(0.45),
                            font_size=16, bold=True, color=colors.get('success', '#276749'))
                add_textbox(slide, html.unescape(desc), Inches(x + 0.15), Inches(1.8), Inches(3.5), Inches(1.2),
                            font_size=13, color=colors.get('text', '#222222'))
        
        # Bullet list
        bullet_items = re.findall(r'bullet-list.*?<li>(.*?)</li>', content, re.DOTALL)
        if bullet_items:
            for i, item in enumerate(bullet_items[:5]):
                text = html.unescape(re.sub(r'<[^>]+>', '', item))
                add_textbox(slide, '— ' + text, Inches(6.5), Inches(1.2 + i * 0.85), Inches(6), Inches(0.45),
                            font_size=16, color=colors.get('text', '#222222'))
        
        # Arch boxes
        arch_boxes = re.findall(r'arch-box.*?<h4>(.*?)</h4>.*?<p>(.*?)</p>', content, re.DOTALL)
        if arch_boxes and len(arch_boxes) >= 4:
            positions = [(1.5, 1.8), (5.0, 1.8), (8.5, 1.8), (4.0, 4.3)]
            for i, (title, desc) in enumerate(arch_boxes[:4]):
                x, y = positions[i]
                w = 3.0 if i < 3 else 4.0
                add_rect(slide, Inches(x), Inches(y), Inches(w), Inches(0.95), '#f5f5f5', colors.get('text', '#222222'))
                add_textbox(slide, html.unescape(title), Inches(x + 0.1), Inches(y + 0.08), Inches(w - 0.2), Inches(0.4),
                            font_size=16, bold=True, color=colors.get('text', '#222222'), align=PP_ALIGN.CENTER)
                add_textbox(slide, html.unescape(desc), Inches(x + 0.1), Inches(y + 0.5), Inches(w - 0.2), Inches(0.4),
                            font_size=12, color=colors.get('text_light', '#555555'), align=PP_ALIGN.CENTER)
        
        # DB tables
        db_headers = re.findall(r'db-table-header">(.*?)</div>', content)
        if db_headers and len(db_headers) >= 5:
            cols_data = [[], [], []]
            for i, header in enumerate(db_headers[:8]):
                cols_data[i % 3].append(html.unescape(header))
            
            for col_idx, tables in enumerate(cols_data):
                base_x = 0.4 + col_idx * 4.2
                for i, table_name in enumerate(tables):
                    y = 1.2 + i * 2.3
                    add_rect(slide, Inches(base_x), Inches(y), Inches(3.8), Inches(2.0),
                             '#f5f5f5', colors.get('accent', '#2c5282'))
                    add_rect(slide, Inches(base_x), Inches(y), Inches(3.8), Inches(0.4),
                             colors.get('accent', '#2c5282'))
                    add_textbox(slide, table_name, Inches(base_x + 0.1), Inches(y + 0.05), Inches(3.6), Inches(0.35),
                                font_size=14, bold=True, color='#ffffff')
                    fields_match = re.search(r'db-table-header">' + re.escape(table_name) + r'</div>.*?<div class="db-table-body">(.*?)</div>', content, re.DOTALL)
                    if fields_match:
                        field_divs = re.findall(r'<div class="db-field[^>]*>(.*?)</div>', fields_match.group(1), re.DOTALL)
                        for fi, field_div in enumerate(field_divs[:5]):
                            ftext = html.unescape(re.sub(r'<[^>]+>', ' ', field_div)).strip()
                            if ftext:
                                add_textbox(slide, ftext, Inches(base_x + 0.15), Inches(y + 0.5 + fi * 0.28), Inches(3.5), Inches(0.28),
                                            font_size=12, color=colors.get('text_light', '#555555'))

        # Screenshot in demo slide (login image)
        demo_img = re.search(r'<img src="figures/([^"]+)"[^>]*alt="([^"]*)"', content)
        if demo_img and 'login' in demo_img.group(1).lower():
            img_path = '/Users/amir/Desktop/back_work_final/figures/' + demo_img.group(1)
            try:
                slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.2), Inches(5.5), Inches(3.5))
            except:
                pass

print(f"\nCreated presentation with {len(prs.slides)} slides")
prs.save('/Users/amir/Desktop/back_work_final/presentation.pptx')
print("Saved to /Users/amir/Desktop/back_work_final/presentation.pptx")