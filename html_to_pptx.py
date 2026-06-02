from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
BG_COLOR = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x2C, 0x52, 0x82)
ACCENT_LIGHT = RGBColor(0x31, 0x82, 0xCE)
TEXT_COLOR = RGBColor(0x22, 0x22, 0x22)
TEXT_LIGHT = RGBColor(0x55, 0x55, 0x55)
WARNING = RGBColor(0xC0, 0x56, 0x21)
SUCCESS = RGBColor(0x27, 0x67, 0x49)
BORDER = RGBColor(0xDD, 0xDD, 0xDD)
BG_ALT = RGBColor(0xF5, 0xF5, 0xF5)

def rect(slide, left, top, width, height, fill=None, border=None):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    if fill:
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

def txt(slide, left, top, width, height, text, size=Pt(20), color=None,
        bold=False, italic=False, align=PP_ALIGN.LEFT, font='Georgia'):
    from pptx.util import Inches, Pt
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    if color:
        run.font.color.rgb = color
    return txBox

def section_label(slide, text):
    txt(slide, Inches(10), Inches(0.28), Inches(2.5), Inches(0.4), text,
        size=Pt(12), color=TEXT_LIGHT, font='Arial')

def heading(slide, text):
    txt(slide, PAD, Inches(0.8), Inches(11.3), Inches(0.8), text,
        size=Pt(32), color=ACCENT, font='Georgia')

W = prs.slide_width
H = prs.slide_height
PAD = Inches(0.6)

# ─── SLIDE 1: Title ───
s1 = prs.slides.add_slide(prs.slide_layouts[6])
# Emblem
circle = s1.shapes.add_shape(9, Inches(6.2), Inches(0.7), Inches(0.9), Inches(0.9))
circle.fill.solid(); circle.fill.fore_color.rgb = ACCENT; circle.line.fill.background()
tf = circle.text_frame; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
run = tf.paragraphs[0].add_run(); run.text = 'АГУ'
run.font.size = Pt(22); run.font.bold = True
run.font.color.rgb = RGBColor(255,255,255); run.font.name = 'Arial'
# Title
txt(s1, Inches(0.8), Inches(1.8), Inches(11.7), Inches(1.5),
    'Информационный ресурс для информационной поддержки и методического сопровождения '
    'образовательной деятельности Института точных наук и цифровых технологий АГУ',
    size=Pt(28), color=TEXT_COLOR, align=PP_ALIGN.CENTER)
txt(s1, Inches(0.8), Inches(3.3), Inches(11.7), Inches(0.6),
    'Выпускная квалификационная работа',
    size=Pt(22), color=TEXT_LIGHT, italic=True, align=PP_ALIGN.CENTER)
# Meta
for i, (label, value) in enumerate([('Студент','Кузьминова А.А.'), ('Научный руководитель','доцент В.В. Бучацкая')]):
    cx = Inches(3.5 + i*3.6)
    txt(s1, cx, Inches(4.1), Inches(2.5), Inches(0.4), label, size=Pt(13), color=TEXT_LIGHT, font='Arial')
    txt(s1, cx, Inches(4.45), Inches(2.5), Inches(0.5), value, size=Pt(20), color=TEXT_COLOR, bold=True, align=PP_ALIGN.CENTER, font='Arial')
# Footer
txt(s1, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.9),
    'Направление 02.03.03 — Математическое обеспечение и администрирование информационных систем\n'
    'Адыгейский государственный университет \u00b7 Майкоп 2026',
    size=Pt(15), color=TEXT_LIGHT, align=PP_ALIGN.CENTER, font='Arial')

# ─── SLIDE 2: Problems ───
s2 = prs.slides.add_slide(prs.slide_layouts[6])
section_label(s2, 'ПРОБЛЕМАТИКА')
heading(s2, 'Проблемы существующей системы')
problems = [
    ('Дублирование информации', 'Разрозненные решения: сайт, соцсети, почтовые рассылки, файлы. Отсутствие согласованности данных.'),
    ('Трудности поиска', 'Нет централизованного хранилища учебно-методических материалов. Преподаватели вынуждены искать в разных местах.'),
    ('Отсутствие единого доступа', 'Нет централизованного управления доступом к информационным ресурсам Института.'),
    ('Высокая трудоёмкость', 'Ручная координация данных, отсутствие автоматизации управленческих процессов.'),
]
card_w = Inches(5.8); card_h = Inches(1.4)
positions = [(0, 0), (6.2, 0), (0, 2.0), (6.2, 2.0)]
for pos, (title, body) in zip(positions, problems):
    lx = PAD + Inches(pos[0]); ly = Inches(2.0) + Inches(pos[1])
    r = rect(s2, lx, ly, card_w, card_h, fill=RGBColor(0xFF,0xF5,0xF5), border=WARNING)
    r.line.width = Pt(1.5)
    txt(s2, lx + Inches(0.2), ly + Inches(0.12), card_w - Inches(0.4), Inches(0.4),
        title, size=Pt(18), color=WARNING, bold=True, font='Arial')
    txt(s2, lx + Inches(0.2), ly + Inches(0.55), card_w - Inches(0.4), Inches(0.8),
        body, size=Pt(16), color=TEXT_LIGHT, font='Arial')

# ─── SLIDE 3: Purpose and Tasks ───
s3 = prs.slides.add_slide(prs.slide_layouts[6])
section_label(s3, 'ЦЕЛЬ И ЗАДАЧИ')
heading(s3, 'Цель и задачи ВКР')
txt(s3, PAD, Inches(1.8), Inches(5.8), Inches(0.5), 'Цель работы', size=Pt(20), color=ACCENT, bold=True, font='Arial')
r = rect(s3, PAD, Inches(2.25), Inches(5.8), Inches(1.2), fill=BG_ALT, border=BORDER)
txt(s3, PAD + Inches(0.2), Inches(2.35), Inches(5.4), Inches(1.0),
    'Разработка информационного ресурса для комплексной информационной поддержки '
    'и методического сопровождения образовательной деятельности Института точных наук и цифровых технологий',
    size=Pt(16), color=TEXT_COLOR)
txt(s3, Inches(7.2), Inches(1.8), Inches(5.5), Inches(0.5), 'Задачи', size=Pt(20), color=ACCENT, bold=True, font='Arial')
tasks = [
    'Анализ предметной области и информационных потоков',
    'Проектирование архитектуры системы и базы данных',
    'Реализация серверной и клиентской частей',
    'Обеспечение безопасности и разграничение прав доступа',
    'Тестирование разработанной системы',
]
for i, task in enumerate(tasks):
    y = Inches(2.25 + i * 0.78)
    circle = s3.shapes.add_shape(9, Inches(7.2), y, Inches(0.32), Inches(0.32))
    circle.fill.solid(); circle.fill.fore_color.rgb = ACCENT; circle.line.fill.background()
    tf = circle.text_frame; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run(); run.text = str(i+1)
    run.font.size = Pt(14); run.font.bold = True
    run.font.color.rgb = RGBColor(255,255,255); run.font.name = 'Arial'
    txt(s3, Inches(7.6), y + Inches(0.03), Inches(5.0), Inches(0.6),
        task, size=Pt(17), color=TEXT_LIGHT, font='Arial')

# ─── SLIDE 4: Architecture ───
s4 = prs.slides.add_slide(prs.slide_layouts[6])
section_label(s4, 'АРХИТЕКТУРА')
heading(s4, 'Архитектура системы')
boxes = [
    ('Frontend\nReact + Redux + MUI', False),
    ('API Gateway\nREST API', True),
    ('Backend\nNode.js + Express', False),
]
bw = Inches(3.0); bh = Inches(1.1); by = Inches(2.5)
for i, (t, hl) in enumerate(boxes):
    lx = Inches(1.5 + i*4.2)
    r = s4.shapes.add_shape(1, lx, by, bw, bh)
    r.fill.solid()
    r.fill.fore_color.rgb = RGBColor(0xE8,0xF0,0xF8) if hl else BG_ALT
    r.line.color.rgb = ACCENT if hl else TEXT_COLOR; r.line.width = Pt(2)
    tf = r.text_frame; tf.word_wrap = True
    lines = t.split('\n')
    for j, line in enumerate(lines):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER; run = p.add_run(); run.text = line
        run.font.size = Pt(18) if j == 0 else Pt(14)
        run.font.bold = (j == 0)
        run.font.color.rgb = ACCENT if hl else TEXT_COLOR; run.font.name = 'Arial'
    if i < 2:
        arrow = s4.shapes.add_textbox(Inches(lx + bw), by + Inches(0.3), Inches(1.2), Inches(0.5))
        arrow_tf = arrow.text_frame; arrow_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        r2 = arrow_tf.paragraphs[0].add_run(); r2.text = '\u21d4'
        r2.font.size = Pt(32); r2.font.color.rgb = TEXT_LIGHT
txt(s4, Inches(6.1), Inches(3.65), Inches(1), Inches(0.6), '\u2193',
    size=Pt(36), color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
bottom = [('PostgreSQL','Реляционная СУБД'), ('Winston','Логирование')]
for i, (name, desc) in enumerate(bottom):
    lx = Inches(3.0 + i*4.0)
    r = s4.shapes.add_shape(1, lx, Inches(4.3), Inches(3.5), Inches(1.1))
    r.fill.solid(); r.fill.fore_color.rgb = BG_ALT; r.line.color.rgb = TEXT_COLOR; r.line.width = Pt(2)
    tf = r.text_frame; tf.word_wrap = True; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run(); run.text = name
    run.font.size = Pt(20); run.font.bold = True; run.font.name = 'Arial'
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; run2 = p2.add_run(); run2.text = desc
    run2.font.size = Pt(15); run2.font.color.rgb = TEXT_LIGHT; run2.font.name = 'Arial'

# ─── SLIDE 5: Tech Stack ───
s5 = prs.slides.add_slide(prs.slide_layouts[6])
section_label(s5, 'ТЕХНОЛОГИИ')
heading(s5, 'Технологический стек')
txt(s5, PAD, Inches(1.8), Inches(5.8), Inches(0.5), 'Backend', size=Pt(22), color=ACCENT, bold=True, font='Arial')
backend = [
    ('Runtime','Node.js 16+','Асинхронное выполнение I/O-операций'),
    ('Framework','Express.js','REST API, middleware, роутинг'),
    ('ORM','Sequelize','Модели данных, миграции, связи'),
    ('DB','PostgreSQL','Реляционная СУБД'),
]
for i, (badge, name, desc) in enumerate(backend):
    y = Inches(2.35 + i * 0.95)
    r = s5.shapes.add_shape(1, PAD, y, Inches(5.8), Inches(0.85))
    r.fill.solid(); r.fill.fore_color.rgb = BG_ALT; r.line.color.rgb = BORDER
    badge_box = s5.shapes.add_shape(1, PAD + Inches(0.1), y + Inches(0.22), Inches(1.0), Inches(0.4))
    badge_box.fill.solid(); badge_box.fill.fore_color.rgb = ACCENT; badge_box.line.fill.background()
    bf = badge_box.text_frame; bf.paragraphs[0].alignment = PP_ALIGN.CENTER
    br = bf.paragraphs[0].add_run(); br.text = badge; br.font.size = Pt(11); br.font.bold = True
    br.font.color.rgb = RGBColor(255,255,255); br.font.name = 'Arial'
    txt(s5, PAD + Inches(1.2), y + Inches(0.08), Inches(2.8), Inches(0.4), name, size=Pt(17), bold=True, font='Arial')
    txt(s5, PAD + Inches(1.2), y + Inches(0.45), Inches(4.4), Inches(0.4), desc, size=Pt(13), color=TEXT_LIGHT, font='Arial')
txt(s5, Inches(7.2), Inches(1.8), Inches(5.5), Inches(0.5), 'Frontend', size=Pt(22), color=ACCENT, bold=True, font='Arial')
frontend = [
    ('Framework','React 18','UI-компоненты, хуки, контекст'),
    ('State','Redux Toolkit','Централизованное состояние'),
    ('UI','Material UI','Компонентная библиотека'),
    ('Build','Vite','Быстрая сборка и HMR'),
]
for i, (badge, name, desc) in enumerate(frontend):
    y = Inches(2.35 + i * 0.95)
    r = s5.shapes.add_shape(1, Inches(7.2), y, Inches(5.5), Inches(0.85))
    r.fill.solid(); r.fill.fore_color.rgb = BG_ALT; r.line.color.rgb = BORDER
    badge_box = s5.shapes.add_shape(1, Inches(7.3), y + Inches(0.22), Inches(1.0), Inches(0.4))
    badge_box.fill.solid(); badge_box.fill.fore_color.rgb = ACCENT; badge_box.line.fill.background()
    bf = badge_box.text_frame; bf.paragraphs[0].alignment = PP_ALIGN.CENTER
    br = bf.paragraphs[0].add_run(); br.text = badge; br.font.size = Pt(11); br.font.bold = True
    br.font.color.rgb = RGBColor(255,255,255); br.font.name = 'Arial'
    txt(s5, Inches(7.2) + Inches(1.2), y + Inches(0.08), Inches(2.8), Inches(0.4), name, size=Pt(17), bold=True, font='Arial')
    txt(s5, Inches(7.2) + Inches(1.2), y + Inches(0.45), Inches(4.4), Inches(0.4), desc, size=Pt(13), color=TEXT_LIGHT, font='Arial')

# ─── SLIDE 6: Database Schema ───
s6 = prs.slides.add_slide(prs.slide_layouts[6])
section_label(s6, 'БАЗА ДАННЫХ')
heading(s6, 'Структура базы данных')

def db_table(slide, left, top, width, header, fields):
    header_h = Inches(0.38)
    field_h = Inches(0.3)
    total_h = header_h + len(fields) * field_h
    r = slide.shapes.add_shape(1, left, top, width, total_h)
    r.fill.solid(); r.fill.fore_color.rgb = BG_ALT; r.line.color.rgb = ACCENT; r.line.width = Pt(2)
    h = slide.shapes.add_shape(1, left, top, width, header_h)
    h.fill.solid(); h.fill.fore_color.rgb = ACCENT; h.line.fill.background()
    h_tf = h.text_frame; h_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    hr = h_tf.paragraphs[0].add_run(); hr.text = header
    hr.font.size = Pt(15); hr.font.bold = True; hr.font.color.rgb = RGBColor(255,255,255); hr.font.name = 'Arial'
    for j, field in enumerate(fields):
        fy = top + header_h + j * field_h
        ftxt = slide.shapes.add_textbox(left + Inches(0.15), fy, width - Inches(0.3), field_h)
        ftf = ftxt.text_frame; ftf.paragraphs[0].alignment = PP_ALIGN.LEFT
        fr = ftf.paragraphs[0].add_run(); fr.text = field; fr.font.size = Pt(13); fr.font.name = 'Courier New'
        if 'PK' in field: fr.font.color.rgb = WARNING; fr.font.bold = True
        elif 'FK' in field: fr.font.color.rgb = SUCCESS
        else: fr.font.color.rgb = TEXT_COLOR

col_w = Inches(4.0)
cx = PAD
db_table(s6, cx, Inches(1.95), col_w, 'Role', ['PK  id','name','permissions'])
db_table(s6, cx, Inches(3.15), col_w, 'Department', ['PK  id','name','code'])
db_table(s6, cx, Inches(4.35), col_w, 'Course', ['PK  id','name','FK  department_id','description'])
cx2 = Inches(4.6)
db_table(s6, cx2, Inches(1.65), col_w, 'User', ['PK  id','email','password_hash','first_name','last_name','FK  role_id','FK  department_id'])
db_table(s6, cx2, Inches(4.25), col_w, 'Schedule', ['PK  id','FK  course_id','date, time, room'])
cx3 = Inches(9.2)
db_table(s6, cx3, Inches(1.95), col_w, 'Group', ['PK  id','name','FK  course_id','year'])
db_table(s6, cx3, Inches(3.15), col_w, 'Material', ['PK  id','title, file_path','FK  user_id','FK  course_id'])
db_table(s6, cx3, Inches(4.35), col_w, 'News', ['PK  id','title, content','FK  author_id'])

# ─── SLIDE 7: Security ───
s7 = prs.slides.add_slide(prs.slide_layouts[6])
section_label(s7, 'БЕЗОПАСНОСТЬ')
heading(s7, 'Механизмы безопасности')
sec = [
    ('Аутентификация','JWT-токены с refresh rotation. Хеширование паролей алгоритмом bcryptjs.'),
    ('Авторизация','RBAC с ролями и middleware-проверками. Защита маршрутов на уровне API.'),
    ('Защита','Helmet для HTTP-заголовков. Rate-limiting. CORS. Валидация входных данных.'),
]
cw = Inches(3.8); ch = Inches(1.5)
for i, (title, body) in enumerate(sec):
    lx = PAD + i * (cw + Inches(0.25))
    r = s7.shapes.add_shape(1, lx, Inches(1.9), cw, ch)
    r.fill.solid(); r.fill.fore_color.rgb = BG_ALT; r.line.color.rgb = BORDER
    txt(s7, lx + Inches(0.2), Inches(2.0), cw - Inches(0.4), Inches(0.4), title, size=Pt(20), color=ACCENT, bold=True, font='Arial')
    txt(s7, lx + Inches(0.2), Inches(2.45), cw - Inches(0.4), Inches(0.9), body, size=Pt(15), color=TEXT_LIGHT, font='Arial')
stats = [('JWT','Access + Refresh'),('bcrypt','Хеширование'),('RBAC','4 роли'),('100','req/min limit')]
sw = Inches(2.8); sh = Inches(1.3)
for i, (val, lbl) in enumerate(stats):
    lx = PAD + i * (sw + Inches(0.2))
    r = s7.shapes.add_shape(1, lx, Inches(3.6), sw, sh)
    r.fill.solid(); r.fill.fore_color.rgb = BG_ALT; r.line.color.rgb = BORDER
    txt(s7, lx, Inches(3.75), sw, Inches(0.5), val, size=Pt(26), color=ACCENT, bold=True, align=PP_ALIGN.CENTER, font='Arial')
    txt(s7, lx, Inches(4.2), sw, Inches(0.5), lbl, size=Pt(14), color=TEXT_LIGHT, align=PP_ALIGN.CENTER, font='Arial')

# ─── SLIDE 8: User Roles ───
s8 = prs.slides.add_slide(prs.slide_layouts[6])
section_label(s8, 'ПОЛЬЗОВАТЕЛИ')
heading(s8, 'Роли и функционал')
rows = [
    ('Студент','Просмотр расписания, доступ к учебным материалам, просмотр новостей, личный кабинет'),
    ('Преподаватель','Управление материалами курсов, создание объявлений, редактирование курсов, просмотр статистики'),
    ('Методист','Модерирование контента, управление кафедрами, контроль качества учебных материалов'),
    ('Администратор','Полный доступ: управление пользователями, ролями, системными настройками, просмотр логов'),
]
cw0 = Inches(2.2); cw1 = Inches(10.3); th = Inches(0.5); rh = Inches(0.85)
r = s8.shapes.add_shape(1, PAD, Inches(1.85), cw0+cw1, th)
r.fill.solid(); r.fill.fore_color.rgb = ACCENT; r.line.fill.background()
txt(s8, PAD + Inches(0.2), Inches(1.87), cw0, th, 'Роль', size=Pt(16), color=RGBColor(255,255,255), bold=True, font='Arial')
txt(s8, PAD + cw0 + Inches(0.3), Inches(1.87), cw1, th, 'Функционал', size=Pt(16), color=RGBColor(255,255,255), bold=True, font='Arial')
for i, (role, func) in enumerate(rows):
    ry = Inches(2.35 + i * rh)
    r = s8.shapes.add_shape(1, PAD, ry, cw0+cw1, rh)
    r.fill.solid()
    r.fill.fore_color.rgb = BG_ALT if i % 2 == 0 else RGBColor(0xFF,0xFF,0xFF)
    r.line.color.rgb = BORDER
    txt(s8, PAD + Inches(0.2), ry + Inches(0.15), cw0, rh, role, size=Pt(17), color=ACCENT, bold=True, font='Arial')
    txt(s8, PAD + cw0 + Inches(0.3), ry + Inches(0.15), cw1, rh, func, size=Pt(16), color=TEXT_COLOR, font='Arial')

# ─── SLIDE 9: Demo ───
s9 = prs.slides.add_slide(prs.slide_layouts[6])
section_label(s9, 'ДЕМОНСТРАЦИЯ')
heading(s9, 'Демонстрация работы системы')
r = s9.shapes.add_shape(1, PAD, Inches(1.9), Inches(5.8), Inches(4.7))
r.fill.solid(); r.fill.fore_color.rgb = BG_ALT; r.line.color.rgb = BORDER
txt(s9, PAD + Inches(0.12), Inches(2.0), Inches(5.6), Inches(0.3), '\u25cf \u25cf \u25cf', size=Pt(14), color=TEXT_LIGHT, font='Arial')
img_path = '/Users/amir/Desktop/back_work_final/figures/01_login.png'
try:
    s9.shapes.add_picture(img_path, PAD + Inches(0.1), Inches(2.3), Inches(5.6), Inches(4.2))
except:
    txt(s9, PAD, Inches(3.8), Inches(5.8), Inches(0.5), '[01_login.png]', size=Pt(14), color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
txt(s9, Inches(7.2), Inches(1.9), Inches(5.5), Inches(0.5), 'Ключевые возможности', size=Pt(22), color=ACCENT, bold=True, font='Arial')
bullets = [
    'Интуитивный интерфейс для всех ролей пользователей',
    'Централизованное хранение учебных материалов',
    'Оперативное обновление расписания',
    'Система уведомлений и объявлений',
    'Администрирование и мониторинг',
]
for i, b in enumerate(bullets):
    txt(s9, Inches(7.2), Inches(2.5 + i*0.65), Inches(5.5), Inches(0.5), '\u2014 ' + b, size=Pt(18), color=TEXT_LIGHT, font='Arial')

# ─── SLIDE 10: Dashboard ───
s10 = prs.slides.add_slide(prs.slide_layouts[6])
section_label(s10, 'ИНТЕРФЕЙС')
heading(s10, 'Панель управления (Dashboard)')
try:
    s10.shapes.add_picture('/Users/amir/Desktop/back_work_final/figures/02_dashboard.png',
                          Inches(3.0), Inches(1.9), Inches(7.0), Inches(4.5))
except:
    txt(s10, Inches(3.0), Inches(3.5), Inches(7.0), Inches(0.5), '[02_dashboard.png]', size=Pt(14), color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
txt(s10, PAD, Inches(6.8), Inches(11.3), Inches(0.5),
    'Централизованная панель со статистикой, быстрыми ссылками и уведомлениями',
    size=Pt(15), color=TEXT_LIGHT, align=PP_ALIGN.CENTER, font='Arial')

# ─── SLIDE 11: Schedule ───
s11 = prs.slides.add_slide(prs.slide_layouts[6])
section_label(s11, 'ИНТЕРФЕЙС')
heading(s11, 'Расписание занятий')
try:
    s11.shapes.add_picture('/Users/amir/Desktop/back_work_final/figures/03_schedule.png',
                          Inches(3.0), Inches(1.9), Inches(7.0), Inches(4.5))
except:
    txt(s11, Inches(3.0), Inches(3.5), Inches(7.0), Inches(0.5), '[03_schedule.png]', size=Pt(14), color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
txt(s11, PAD, Inches(6.8), Inches(11.3), Inches(0.5),
    'Фильтрация по группам, датам и аудиториям',
    size=Pt(15), color=TEXT_LIGHT, align=PP_ALIGN.CENTER, font='Arial')

# ─── SLIDE 12: Materials ───
s12 = prs.slides.add_slide(prs.slide_layouts[6])
section_label(s12, 'ИНТЕРФЕЙС')
heading(s12, 'Учебные материалы')
try:
    s12.shapes.add_picture('/Users/amir/Desktop/back_work_final/figures/04_materials.png',
                          Inches(3.0), Inches(1.9), Inches(7.0), Inches(4.5))
except:
    txt(s12, Inches(3.0), Inches(3.5), Inches(7.0), Inches(0.5), '[04_materials.png]', size=Pt(14), color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
txt(s12, PAD, Inches(6.8), Inches(11.3), Inches(0.5),
    'Каталог с поиском и категоризацией по дисциплинам',
    size=Pt(15), color=TEXT_LIGHT, align=PP_ALIGN.CENTER, font='Arial')

# ─── SLIDE 13: News ───
s13 = prs.slides.add_slide(prs.slide_layouts[6])
section_label(s13, 'ИНТЕРФЕЙС')
heading(s13, 'Новости и объявления')
try:
    s13.shapes.add_picture('/Users/amir/Desktop/back_work_final/figures/06_news.png',
                          Inches(3.0), Inches(1.9), Inches(7.0), Inches(4.5))
except:
    txt(s13, Inches(3.0), Inches(3.5), Inches(7.0), Inches(0.5), '[06_news.png]', size=Pt(14), color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
txt(s13, PAD, Inches(6.8), Inches(11.3), Inches(0.5),
    'Лента объявлений с возможностью модерирования',
    size=Pt(15), color=TEXT_LIGHT, align=PP_ALIGN.CENTER, font='Arial')

# ─── SLIDE 14: Conclusion ───
s14 = prs.slides.add_slide(prs.slide_layouts[6])
circle = s14.shapes.add_shape(9, Inches(6.1), Inches(2.3), Inches(1.1), Inches(1.1))
circle.fill.solid(); circle.fill.fore_color.rgb = SUCCESS; circle.line.fill.background()
tf = circle.text_frame; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
run = tf.paragraphs[0].add_run(); run.text = '\u2713'
run.font.size = Pt(44); run.font.bold = True; run.font.color.rgb = RGBColor(255,255,255)
txt(s14, Inches(0.8), Inches(3.6), Inches(11.7), Inches(1.2),
    'Спасибо за внимание', size=Pt(36), color=TEXT_COLOR, align=PP_ALIGN.CENTER, font='Georgia')

out = '/Users/amir/Desktop/back_work_final/presentation.pptx'
prs.save(out)
print(f'Saved: {out}')
